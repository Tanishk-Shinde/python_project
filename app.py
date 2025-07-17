from flask import Flask, request, render_template, redirect, url_for, session
from werkzeug.security import generate_password_hash, check_password_hash
import os, json, fitz, sqlite3
from pptx import Presentation
from datetime import datetime
from textblob import TextBlob
import textstat
import whisper

app = Flask(__name__)
app.secret_key = 'secret_key'
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# === DATABASE SETUP ===
def init_db():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL
        )
    ''')
    conn.commit()
    conn.close()

init_db()

def create_user(username, password):
    hashed = generate_password_hash(password)
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, hashed))
    conn.commit()
    conn.close()

def verify_user(username, password):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute("SELECT id, password FROM users WHERE username = ?", (username,))
    user = c.fetchone()
    conn.close()
    if user and check_password_hash(user[1], password):
        return user[0]
    return None

# === TEXT & AUDIO PROCESSING ===
def extract_text_from_pdf(path):
    text = ""
    doc = fitz.open(path)
    for page in doc:
        text += page.get_text()
    doc.close()
    return text

def extract_text_from_pptx(path):
    text = ""
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def transcribe_audio(path):
    model = whisper.load_model("base")
    result = model.transcribe(path)
    return result["text"]

# === NLP / AI ANALYSIS ===
SECTION_KEYWORDS = {
    "Problem": ["problem", "pain", "issue", "need"],
    "Solution": ["solution", "approach", "how we solve"],
    "Market": ["market", "audience", "target"],
    "Business Model": ["model", "revenue", "monetize"],
    "Competition": ["competition", "rival", "differentiation"],
    "Team": ["team", "founder", "background"],
    "Traction": ["traction", "growth", "milestone"],
    "Financials": ["finance", "cost", "funding", "projection"]
}

def detect_sections(text):
    sections = {}
    for section, keywords in SECTION_KEYWORDS.items():
        found = any(kw in text.lower() for kw in keywords)
        sections[section] = "✅" if found else "❌"
    return sections

def score_clarity(text):
    grade = textstat.flesch_kincaid_grade(text)
    return round(10 - min(grade, 10), 2)

def analyze_tone(text):
    polarity = TextBlob(text).sentiment.polarity
    if polarity > 0.3: return "Positive"
    elif polarity < -0.3: return "Negative"
    return "Neutral"

def score_completeness(sections):
    found = sum(1 for v in sections.values() if v == "✅")
    return round((found / len(sections)) * 10, 2)

def analyze_pitch(text):
    sections = detect_sections(text)
    clarity = score_clarity(text)
    tone = analyze_tone(text)
    completeness = score_completeness(sections)
    return {
        "Sections": sections,
        "Clarity Score": clarity,
        "Tone": tone,
        "Completeness Score": completeness
    }

def save_result(analysis, user_id):
    path = 'results.json'
    try:
        with open(path, 'r') as f:
            data = json.load(f)
    except:
        data = []

    data.append({
        "id": len(data) + 1,
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "clarity": analysis["Clarity Score"],
        "completeness": analysis["Completeness Score"],
        "tone": analysis["Tone"],
        "user_id": user_id
    })

    with open(path, 'w') as f:
        json.dump(data, f, indent=2)

@app.route('/')
def root():
    if 'user_id' in session:
        return redirect('/dashboard')
    return redirect('/home')

@app.route('/home')
def home():
    return render_template('home.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        create_user(username, password)
        return redirect('/success')
    return render_template('register.html')

@app.route('/success')
def success():
    return render_template('success.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        user_id = verify_user(username, password)
        if user_id:
            session['username'] = username
            session['user_id'] = user_id
            return redirect('/upload')
        return 'Invalid credentials. Try again.'
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect('/login')

@app.route('/upload', methods=['GET', 'POST'])
def upload():
    if 'user_id' not in session:
        return redirect('/login')

    if request.method == 'POST':
        pitch_file = request.files.get('pitch_file')
        audio_file = request.files.get('audio_file')
        text = ""

        if pitch_file:
            pitch_path = os.path.join(UPLOAD_FOLDER, pitch_file.filename)
            pitch_file.save(pitch_path)
            if pitch_path.endswith('.pdf'):
                text = extract_text_from_pdf(pitch_path)
            elif pitch_path.endswith('.pptx'):
                text = extract_text_from_pptx(pitch_path)

        if audio_file:
            audio_path = os.path.join(UPLOAD_FOLDER, audio_file.filename)
            audio_file.save(audio_path)
            transcript = transcribe_audio(audio_path)
            text += "\n" + transcript

        analysis = analyze_pitch(text)
        save_result(analysis, session.get('user_id'))
        return redirect(url_for('dashboard'))

    return render_template('upload.html', username=session.get("username", "User"))

@app.route('/dashboard')
def dashboard():
    if 'user_id' not in session:
        return redirect('/login')

    results = []
    try:
        with open('results.json', 'r') as f:
            all_results = json.load(f)
            results = [r for r in all_results if str(r.get('user_id')) == str(session['user_id'])]
    except (FileNotFoundError, json.JSONDecodeError) as e:
        print(f"Error loading results: {e}")

    return render_template('dashboard.html', results=results)

if __name__ == '__main__':
    app.run(debug=True)
