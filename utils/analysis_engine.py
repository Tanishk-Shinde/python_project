import fitz, whisper, torch
from pptx import Presentation
from textblob import TextBlob
from textstat import flesch_reading_ease

def extract_text(path):
    if path.endswith('.pdf'):
        doc = fitz.open(path)
        return "\n".join([page.get_text() for page in doc])
    elif path.endswith('.ppt') or path.endswith('.pptx'):
        prs = Presentation(path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return ""

def analyze_pitch(filepath):
    if filepath.endswith(('.pdf', '.ppt', '.pptx')):
        text = extract_text(filepath)
        blob = TextBlob(text)
        grammar_score = 10 - len(blob.correct().split()) / len(blob.words)
        readability = flesch_reading_ease(text)
        return {
            "type": "text",
            "word_count": len(text.split()),
            "grammar_score": round(grammar_score * 10, 2),
            "readability_score": round(readability, 2),
            "final_score": round((grammar_score * 10 + readability) / 2, 2)
        }
    elif filepath.endswith('.wav') or filepath.endswith('.mp3'):
        model = whisper.load_model("base")
        result = model.transcribe(filepath)
        text = result["text"]
        blob = TextBlob(text)
        grammar_score = 10 - len(blob.correct().split()) / len(blob.words)
        return {
            "type": "audio",
            "transcription": text,
            "grammar_score": round(grammar_score * 10, 2),
            "final_score": round(grammar_score * 10, 2)
        }
    else:
        return {"error": "Unsupported file type"}
